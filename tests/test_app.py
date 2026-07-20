from pathlib import Path
from concurrent.futures import ThreadPoolExecutor
from unittest.mock import patch
from types import SimpleNamespace
import asyncio
from datetime import datetime
from decimal import Decimal
import hashlib
import hmac
import json
import time

import pymupdf
import pytest

from fastapi.testclient import TestClient

import app as application
from services.answer_service import AnswerService, NOT_FOUND_MESSAGE, format_abnt_references, format_sources
from services.auth_repository import AuthRepository
from services.vector_store import LocalVectorStore
from services.query_analysis import QueryType, analyze_query
from services.rag_diagnostics import RAGDiagnosticsRepository
from services.document_loader import load_document
from services.editorial_style import HOMILY_CORPUS_PROFILE
from services.presentation_service import PresentationService, safe_filename
from services.localization import answer_message
from services.asaas_service import AsaasService
from services.mercado_pago_service import MercadoPagoError, MercadoPagoService
from services.subscription_service import SubscriptionService


def authenticated_client() -> TestClient:
    client = TestClient(application.app)
    if not application.auth_repository.find_user_by_login("Admin"):
        application.auth_repository.ensure_admin("AdminTest3510")
    for password in ("AdminTest3510", "3510"):
        response = client.post("/login", data={"email": "Admin", "senha": password})
        if response.url.path == "/":
            return client
    raise AssertionError("Nao foi possivel autenticar o administrador de teste.")


def create_free_user(client: TestClient, email: str = "teste@exemplo.com", password: str = "Senha123") -> None:
    response = client.post("/cadastro", data={"nome": "Teste Usuario", "email": email, "senha": password}, follow_redirects=False)
    assert response.status_code == 303
    login = client.post("/login", data={"email": email, "senha": password})
    assert login.status_code == 200


def test_vector_store_indexes_and_finds_txt(tmp_path: Path):
    documents = tmp_path / "Documentos"
    documents.mkdir()
    (documents / "catequese.txt").write_text(
        "A esperança cristã está fundada na ressurreição de Jesus Cristo.", encoding="utf-8"
    )
    store = LocalVectorStore(documents, tmp_path / "index" / "indice.json", 300, 40)
    progress = []
    status = store.index_documents(lambda current, total, name: progress.append((current, total, name)))
    results = store.search("Em que está fundada a esperança cristã?", minimum_score=0.01)

    assert status["documentos"] == 1
    assert status["trechos"] == 1
    assert results[0]["source"] == "catequese.txt"
    assert progress[-1] == (1, 1, "catequese.txt")

    with patch("services.vector_store.load_document", side_effect=AssertionError("não deveria reler")):
        repeated_status = store.index_documents()

    assert repeated_status["documentos"] == 1
    assert repeated_status["trechos"] == 1


def test_pdf_fast_reader(tmp_path: Path):
    documents = tmp_path / "Documentos"
    documents.mkdir()
    pdf_path = documents / "fonte.pdf"
    with pymupdf.open() as pdf:
        page = pdf.new_page()
        page.insert_text((72, 72), "Documento pastoral de teste")
        pdf.save(pdf_path)

    sections = list(load_document(pdf_path, documents))

    assert sections[0].source == "fonte.pdf"
    assert sections[0].location == "página 1"
    assert "Documento pastoral" in sections[0].text
def test_explicit_source_name_filters_results(tmp_path: Path):
    documents = tmp_path / "Documentos"
    documents.mkdir()
    (documents / "Catecismo.txt").write_text("Os sacramentos comunicam a graça de Cristo.", encoding="utf-8")
    (documents / "Outro livro.txt").write_text("Os sacramentos são sinais importantes.", encoding="utf-8")
    store = LocalVectorStore(documents, tmp_path / "indice.json", 300, 40)
    store.index_documents()

    results = store.search("O que o Catecismo ensina sobre os sacramentos?", minimum_score=0)

    assert results
    assert {item["source"] for item in results} == {"Catecismo.txt"}


def test_editorial_source_order(tmp_path: Path):
    documents = tmp_path / "Documentos"
    documents.mkdir()
    fixtures = {
        "Catecismo da Igreja Católica.txt": "A caridade nasce do amor de Deus.",
        "compendio-dos-simbolos-definicoes-e-declaracoes-de-fe-e-moral.txt": "A caridade é uma virtude.",
        "compendio-da-doutrina-social-da-igreja.txt": "A caridade orienta a vida social.",
        "Suma Teológica.txt": "A caridade une a pessoa a Deus.",
        "Bíblia Ave Maria - Edição de Estudo.txt": "A caridade é paciente.",
        "Compêndio Vaticano II.txt": "A caridade manifesta a comunhão da Igreja.",
        "A Fé Explicada.txt": "A caridade é uma virtude explicada.",
        "Documento complementar.txt": "A caridade transforma a comunidade.",
    }
    for name, text in fixtures.items():
        (documents / name).write_text(text, encoding="utf-8")
    store = LocalVectorStore(documents, tmp_path / "indice.json", 300, 40)
    store.index_documents()

    results = store.search_ordered("O que é a caridade?", minimum_score=0)

    categories = list(dict.fromkeys(item["categoria"] for item in results))
    assert categories[:7] == [
        "Catecismo da Igreja Católica",
        "Bíblia Ave Maria — citações bíblicas",
        "Compêndio Vaticano II",
        "Compêndio da Doutrina Social da Igreja",
        "A Fé Explicada",
        "Compêndio dos símbolos, definições e declarações",
        "Suma Teológica",
    ]
    assert categories[-1] == "Demais documentos"


def test_single_word_uses_nominal_index_hierarchy(tmp_path: Path):
    documents = tmp_path / "Documentos"
    documents.mkdir()
    fixtures = {
        "Catecismo da Igreja Católica.txt": (
            "SUMÁRIO\nPECADO ................................ 1846-1876\n\n"
            "1846. A misericórdia divina vence toda transgressão humana."
        ),
        "compendio-dos-simbolos-definicoes-e-declaracoes-de-fe-e-moral.txt": (
            "ÍNDICE SISTEMÁTICO\nPecado: natureza D:1c\n\n"
            "D:1c A desobediência moral rompe a comunhão com Deus."
        ),
        "A Fé Explicada.txt": (
            "SUMÁRIO\nO que é o pecado? ........................ 76\n\n"
            "[p. 76] A recusa consciente do amor divino fere a vida espiritual."
        ),
        "Bíblia Ave Maria.txt": "O pecado entrou no mundo, mas a graça foi abundante.",
    }
    for name, text in fixtures.items():
        (documents / name).write_text(text, encoding="utf-8")
    store = LocalVectorStore(documents, tmp_path / "indice.json", 105, 5)
    store.index_documents()

    results = store.search_ordered("PECADO", minimum_score=0)

    categories = list(dict.fromkeys(item["categoria"] for item in results))
    assert categories[:3] == [
        "Catecismo da Igreja Católica",
        "Compêndio dos símbolos, definições e declarações",
        "A Fé Explicada",
    ]
    assert any("transgressão humana" in item["text"] for item in results)
    assert any("desobediência moral" in item["text"] for item in results)
    assert any("recusa consciente" in item["text"] for item in results)


def test_definition_question_prioritizes_the_direct_definition(tmp_path: Path):
    documents = tmp_path / "Documentos"
    documents.mkdir()
    (documents / "Catecismo da Igreja Católica.txt").write_text(
        "A remissão dos pecados é concedida no Batismo.\n\n"
        "II. A definição do pecado\n1849. O pecado é uma falta contra a razão, "
        "a verdade e a consciência reta.",
        encoding="utf-8",
    )
    (documents / "A Fé Explicada.txt").write_text(
        "SUMÁRIO\nO PECADO ATUAL ........ 71\n\n"
        "## P?gina 71\nO PECADO ATUAL\nO pecado atual é aquele que nós mesmos cometemos.",
        encoding="utf-8",
    )
    store = LocalVectorStore(documents, tmp_path / "indice.json", 180, 10)
    store.index_documents()

    for query in ("PECADO", "o que é pecado"):
        results = store.search_ordered(query, minimum_score=0)
        catechism = [item for item in results if "Catecismo" in item["source"]]
        assert catechism
        assert "definição do pecado" in catechism[0]["text"]
        assert any("nós mesmos cometemos" in item["text"] for item in results)


def test_single_word_searches_remaining_documents_before_negative(tmp_path: Path):
    documents = tmp_path / "Documentos"
    documents.mkdir()
    (documents / "Documento complementar.txt").write_text(
        "A perseverança sustenta a esperança nas provações.",
        encoding="utf-8",
    )
    store = LocalVectorStore(documents, tmp_path / "indice.json", 300, 40)
    store.index_documents()

    results = store.search_ordered("PERSEVERANÇA", minimum_score=0.08)

    assert results
    assert results[0]["categoria"] == "Demais documentos"


@pytest.mark.parametrize(
    ("query", "expected"),
    [
        ("PECADO", QueryType.TERM),
        ("pecado mortal", QueryType.PHRASE),
        ("O que é pecado mortal?", QueryType.QUESTION),
        ("CIC 1854", QueryType.REFERENCE),
        ("Explique a graça", QueryType.COMMAND),
    ],
)
def test_query_classification_never_blocks_retrieval(query: str, expected: QueryType):
    analysis = analyze_query(query)

    assert analysis.query_type == expected
    assert analysis.lexical_terms
    assert analysis.expanded_queries[0] == analysis.normalized


def test_versioned_catholic_single_term_fixture_is_complete_and_classified():
    fixture = json.loads(
        (Path(__file__).parent / "fixtures" / "catholic_single_term_queries.json").read_text(encoding="utf-8")
    )

    assert fixture["version"] == 1
    assert fixture["language"] == "pt-BR"
    assert len(fixture["queries"]) >= 500
    for item in fixture["queries"]:
        analysis = analyze_query(item["term"])
        assert analysis.query_type == QueryType.TERM, item["term"]
        assert analysis.lexical_terms, item["term"]
        assert analysis.expanded_queries[0] == item["term"]


def test_required_rag_regression_queries_use_all_retrieval_paths(tmp_path: Path):
    documents = tmp_path / "Documentos"
    documents.mkdir()
    terms = [
        "pecado", "graça", "Maria", "Eucaristia", "purgatório", "confissão",
        "Trindade", "indulgência", "batismo", "justiça", "caridade", "escatologia",
        "liturgia", "ressurreição",
    ]
    (documents / "Catecismo temático.txt").write_text(
        "\n\n".join(f"{term}. Seção documental sobre {term}." for term in terms),
        encoding="utf-8",
    )
    store = LocalVectorStore(documents, tmp_path / "indice.json", 120, 10)
    store.index_documents()

    for term in terms:
        results, diagnostics = store.search_ordered(term.upper(), include_diagnostics=True)
        assert results, term
        assert diagnostics["query"]["query_type"] == "TERM"
        assert diagnostics["candidate_counts"]["lexical_exact"] > 0
        assert diagnostics["candidate_counts"]["global_fallback"] > 0


def test_two_character_term_and_accent_variants_are_valid(tmp_path: Path):
    documents = tmp_path / "Documentos"
    documents.mkdir()
    (documents / "Catecismo.txt").write_text(
        "A fé é a resposta do ser humano a Deus. A bênção é dom e encontro.",
        encoding="utf-8",
    )
    store = LocalVectorStore(documents, tmp_path / "indice.json", 300, 40)
    store.index_documents()

    faith = store.search_ordered("Fé", minimum_score=0.08)
    accented = store.search_ordered("BÊNÇÃO", minimum_score=0.08)
    unaccented = store.search_ordered("bencao", minimum_score=0.08)

    assert faith
    assert [item["id"] for item in accented] == [item["id"] for item in unaccented]


def test_exact_lexical_match_survives_high_requested_threshold(tmp_path: Path):
    documents = tmp_path / "Documentos"
    documents.mkdir()
    (documents / "Documento complementar.txt").write_text(
        "Escatologia é o tema desta seção documental.", encoding="utf-8"
    )
    store = LocalVectorStore(documents, tmp_path / "indice.json", 300, 40)
    store.index_documents()

    results, diagnostics = store.search_ordered(
        "ESCATOLOGIA", minimum_score=99, include_diagnostics=True
    )

    assert results
    assert "lexical_exact" in results[0]["retrieval_strategies"]
    assert diagnostics["threshold_policy"] == "dynamic_exact_and_lexical_matches_are_protected"


def test_equivalent_case_queries_reuse_safe_search_cache(tmp_path: Path):
    documents = tmp_path / "Documentos"
    documents.mkdir()
    (documents / "Catecismo.txt").write_text("O pecado fere a comunhão.", encoding="utf-8")
    store = LocalVectorStore(documents, tmp_path / "indice.json", 300, 40)
    store.index_documents()

    with patch.object(store, "_exact_candidate_rows", wraps=store._exact_candidate_rows) as exact_search:
        uppercase = store.search_ordered("PECADO", include_diagnostics=True)
        lowercase = store.search_ordered("pecado", include_diagnostics=True)

    assert exact_search.call_count == 1
    assert [item["id"] for item in uppercase[0]] == [item["id"] for item in lowercase[0]]


def test_editorial_bonus_does_not_make_irrelevant_chunk_relevant(tmp_path: Path):
    documents = tmp_path / "Documentos"
    documents.mkdir()
    (documents / "Catecismo da Igreja Católica.txt").write_text(
        "A música litúrgica acompanha a celebração.", encoding="utf-8"
    )
    store = LocalVectorStore(documents, tmp_path / "indice.json", 300, 40)
    store.index_documents()

    results = store.search_ordered(
        "Explique detalhadamente a engenharia de foguetes interplanetários.",
        minimum_score=0.08,
    )

    assert results == []


def test_stream_continues_after_output_limit(monkeypatch):
    class AsyncEvents:
        def __init__(self, events):
            self.events = iter(events)

        def __aiter__(self):
            return self

        async def __anext__(self):
            try:
                return next(self.events)
            except StopIteration as exc:
                raise StopAsyncIteration from exc

    class FakeResponses:
        def __init__(self):
            self.calls = 0

        async def create(self, **kwargs):
            self.calls += 1
            response_id = f"resposta-{self.calls}"
            events = [
                SimpleNamespace(type="response.created", response=SimpleNamespace(id=response_id)),
                SimpleNamespace(type="response.output_text.delta", delta="Texto interrompido" if self.calls == 1 else " e concluído."),
                SimpleNamespace(
                    type="response.incomplete" if self.calls == 1 else "response.completed",
                    response=SimpleNamespace(id=response_id),
                ),
            ]
            return AsyncEvents(events)

    fake_responses = FakeResponses()
    monkeypatch.setattr(
        "services.answer_service.AsyncOpenAI",
        lambda api_key: SimpleNamespace(responses=fake_responses),
    )
    service = AnswerService("chave", "modelo")
    chunks = [{"source": "Catecismo.txt", "location": "página 1", "text": "Trecho", "score": 1.0}]

    async def collect():
        return "".join([part async for part in service.stream_answer("Pergunta", chunks)])

    answer = asyncio.run(collect())

    assert answer == "Texto interrompido e concluído."
    assert fake_responses.calls == 2


def test_answer_is_rewritten_when_reviewer_rejects(monkeypatch):
    class FakeResponses:
        def __init__(self):
            self.calls = []

        async def create(self, **kwargs):
            self.calls.append(kwargs)
            if len(self.calls) == 1:
                return SimpleNamespace(output_text="Resposta original com extrapolação.")
            if len(self.calls) == 2:
                return SimpleNamespace(
                    output_text='{"action": "block", "reason": "Extrapolou a base", "suggested_answer": "Não encontrei conteúdo correspondente na base documental."}'
                )
            return SimpleNamespace(output_text="Resposta reescrita apenas com o trecho cadastrado.")

    fake_responses = FakeResponses()
    monkeypatch.setattr(
        "services.answer_service.AsyncOpenAI",
        lambda api_key: SimpleNamespace(responses=fake_responses),
    )
    service = AnswerService("chave", "modelo", "modelo-revisor")
    chunks = [{"source": "Catecismo.txt", "location": "página 1", "text": "Trecho", "score": 1.0}]

    async def call():
        return await service.answer_with_review("Pergunta", chunks)

    result = asyncio.run(call())

    assert result["resposta"] == "Resposta reescrita apenas com o trecho cadastrado."
    assert result["status_revisao"] == "rewrite"
    assert result["motivo_revisao"] == "Extrapolou a base"
    assert len(fake_responses.calls) == 3


def test_answer_softens_generic_block_to_rewrite(monkeypatch):
    class FakeResponses:
        def __init__(self):
            self.calls = []

        async def create(self, **kwargs):
            self.calls.append(kwargs)
            if len(self.calls) == 1:
                return SimpleNamespace(output_text="Resposta simples e apoiada.")
            return SimpleNamespace(
                output_text='{"action": "block", "reason": "Não foi possível validar a resposta com segurança.", "suggested_answer": "Não encontrei essa informação nos documentos cadastrados."}'
            )

    fake_responses = FakeResponses()
    monkeypatch.setattr(
        "services.answer_service.AsyncOpenAI",
        lambda api_key: SimpleNamespace(responses=fake_responses),
    )
    service = AnswerService("chave", "modelo", "modelo-revisor")
    chunks = [{"source": "Catecismo.txt", "location": "página 1", "text": "Trecho", "score": 1.0}]

    async def call():
        return await service.answer_with_review("Pergunta", chunks)

    result = asyncio.run(call())

    assert result["status_revisao"] == "rewrite"
    assert result["resposta"] == "Resposta simples e apoiada."


def test_answer_prompt_requires_consolidated_text():
    service = AnswerService("chave", "modelo")
    chunks = [{"source": "Catecismo.txt", "location": "página 1", "text": "Trecho", "score": 1.0}]

    request = service._request_arguments(
        "Pergunta",
        chunks,
        [],
        [{"source": "joao-paulo-ii-homilias-txt/homilia.txt", "location": "documento", "text": "Amados irmaos."}],
    )
    instructions = request["instructions"]

    assert "uma única síntese consolidada" in instructions
    assert "Não informe nem liste as fontes no corpo" in instructions
    assert "A Fé Explicada" in instructions
    assert "PADRÃO HOMILÉTICO DE SÃO JOÃO PAULO II" in instructions
    assert "português brasileiro contemporâneo" in instructions
    assert "O padrão rege somente a forma" in instructions
    assert "AMOSTRAS DE ESTILO DAS HOMILIAS" in instructions
    assert "AMOSTRAS DE ESTILO DAS HOMILIAS" in request["input"]


def test_catechesis_prompt_goes_directly_to_topic_and_adapts_examples_to_audience():
    service = AnswerService("chave", "modelo")
    chunks = [{"source": "Catecismo.txt", "location": "página 1", "text": "Trecho", "score": 1.0}]

    request = service._request_arguments(
        "Redija uma catequese sobre a esperança cristã para crianças de 8 anos, com exemplos didáticos.",
        chunks,
        [],
    )
    instructions = request["instructions"]

    assert "não defina, explique nem introduza o que é uma catequese" in instructions
    assert "exemplos concretos adequados ao público declarado" in instructions
    assert service._is_catechesis_request("Prepare uma catequese sobre o perdão para adolescentes.")
    assert service._is_catechesis_request("Elabore uma catequese sobre a Eucaristia.")
    assert not service._is_catechesis_request("O que é uma catequese?")


def test_homily_corpus_profile_and_presentation_standard():
    instructions = PresentationService._plan_instructions()

    assert HOMILY_CORPUS_PROFILE["documents"] == 1093
    assert HOMILY_CORPUS_PROFILE["period"] == "1978-2005"
    assert HOMILY_CORPUS_PROFILE["words"] == 1_662_755
    assert "PADRÃO HOMILÉTICO DE SÃO JOÃO PAULO II" in instructions
    assert "anúncio do tema, aprofundamento" in instructions
    assert "títulos como afirmações vivas e sóbrias" in instructions
    assert "sem acrescentar conteúdo novo" in instructions


def test_free_plan_limits_queries_and_presentations(tmp_path: Path, monkeypatch):
    db_file = tmp_path / "magisteria.sqlite"
    repo = AuthRepository(db_file)
    monkeypatch.setattr(application, "auth_repository", repo)
    client = TestClient(application.app)
    create_free_user(client, email="livre@exemplo.com")

    class FakeAnswerService:
        api_key = "x"

        async def answer_with_review(self, *args, **kwargs):
            return {"resposta": "ok", "status_revisao": "approve", "motivo_revisao": ""}

    monkeypatch.setattr(application, "answer_service", FakeAnswerService())
    monkeypatch.setattr(application, "ordered_chunks", lambda payload: [{"source": "x", "location": "página 1", "text": "trecho", "score": 1.0}])
    monkeypatch.setattr(application, "homily_style_chunks", lambda payload: [])
    monkeypatch.setattr(application.vector_store, "search_ordered", lambda *args, **kwargs: [{"source": "x", "location": "página 1", "text": "trecho", "ordem": 1, "categoria": "Documento"}])
    async def fake_create_outline(*args, **kwargs):
        return [{"titulo": "A", "sintese": "B", "pontos": ["c"]}]

    async def fake_create_plan(*args, **kwargs):
        return {"titulo_curto": "A", "frase_final": "B", "topicos": [{"titulo": "A", "sintese": "B", "pontos": ["c"]}]}

    async def fake_create_pptx(*args, **kwargs):
        return b"pptx"

    monkeypatch.setattr(application.presentation_service, "create_outline", fake_create_outline)
    monkeypatch.setattr(application.presentation_service, "create_docx", lambda *args, **kwargs: b"docx")
    monkeypatch.setattr(application.presentation_service, "create_plan", fake_create_plan)
    monkeypatch.setattr(application.presentation_service, "create_pptx", fake_create_pptx)

    for _ in range(3):
        response = client.post("/perguntar-stream", json={"pergunta": "teste", "historico": []})
        assert response.status_code == 200

    blocked = client.post("/perguntar-stream", json={"pergunta": "teste", "historico": []})
    assert blocked.status_code == 403

    script = client.post("/criar-roteiro", json={"titulo": "Tema", "resposta": "Texto suficiente para gerar roteiro"})
    assert script.status_code == 200
    slide = client.post("/criar-slides", json={"titulo": "Tema", "resposta": "Texto suficiente para gerar slides"})
    assert slide.status_code == 200

    blocked_script = client.post("/criar-roteiro", json={"titulo": "Tema", "resposta": "Texto suficiente para gerar roteiro"})
    blocked_slide = client.post("/criar-slides", json={"titulo": "Tema", "resposta": "Texto suficiente para gerar slides"})
    assert blocked_script.status_code == 403
    assert blocked_slide.status_code == 403


def test_coupon_and_verified_payment_activate_full_access(tmp_path: Path, monkeypatch):
    db_file = tmp_path / "magisteria.sqlite"
    repo = AuthRepository(db_file)
    monkeypatch.setattr(application, "auth_repository", repo)
    monkeypatch.setattr(application, "FREE_CUPON_CODES", {"LIBERA100"})

    class FakeMercadoPago:
        configured = True
        webhook_signature_configured = True
        price = Decimal("29.90")
        currency = "BRL"
        invoice_status = "approved"

        async def create_subscription(self, user, external_reference):
            self.reference = external_reference
            return {
                "id": "sub_123",
                "checkout_url": "https://www.mercadopago.com.br/checkout/v1/redirect",
            }

        def validate_webhook_signature(self, *args):
            return True

        async def get_subscription(self, subscription_id):
            return {
                "id": subscription_id,
                "external_reference": self.reference,
                "status": "authorized",
                "auto_recurring": {"transaction_amount": 29.9, "currency_id": "BRL"},
                "next_payment_date": "2026-08-14T12:00:00Z",
            }

        async def get_authorized_payment(self, invoice_id):
            return {
                "id": invoice_id,
                "preapproval_id": "sub_123",
                "external_reference": self.reference,
                "transaction_amount": 29.9,
                "currency_id": "BRL",
                "payment": {"id": "123456", "status": self.invoice_status, "status_detail": "accredited"},
            }

    provider = FakeMercadoPago()
    monkeypatch.setattr(application, "mercado_pago_service", provider)
    client = TestClient(application.app)
    create_free_user(client, email="cupom@exemplo.com")

    coupon = client.post("/assinatura/cupom", json={"cupom": "LIBERA100"})
    assert coupon.status_code == 200
    assert coupon.json()["usuario"]["is_full_access"] is True

    other = TestClient(application.app)
    create_free_user(other, email="mp@exemplo.com")
    mp_user = repo.find_user_by_login("mp@exemplo.com")
    order = repo.create_payment_order(mp_user["id"], "29.90", "BRL", "mercado_pago")
    repo.attach_payment_preference(order["reference"], "sub_123")
    provider.reference = order["reference"]
    order = repo.get_payment_order(order["reference"])
    assert order["provider_preference_id"] == "sub_123"

    webhook = TestClient(application.app).post(
        "/webhooks/mercadopago?data.id=sub_123&type=subscription_preapproval",
        headers={"x-signature": "ts=1,v1=fake", "x-request-id": "request-1"},
        json={"type": "subscription_preapproval", "data": {"id": "sub_123"}, "status": "forged"},
    )
    assert webhook.status_code == 200
    assert webhook.json()["status"] == "authorized"
    updated = repo.find_user_by_login("mp@exemplo.com")
    assert updated["account_type"] == "completa"
    assert updated["payment_provider_subscription_id"] == "sub_123"

    provider.invoice_status = "refunded"
    refund = TestClient(application.app).post(
        "/webhooks/mercadopago?data.id=789&type=subscription_authorized_payment",
        headers={"x-signature": "ts=2,v1=fake", "x-request-id": "request-3"},
        json={"type": "subscription_authorized_payment", "data": {"id": "789"}},
    )
    assert refund.status_code == 200
    assert repo.find_user_by_login("mp@exemplo.com")["account_type"] == "gratuita"


def test_admin_creates_managed_coupon_user_redeems_and_admin_revokes(tmp_path: Path, monkeypatch):
    repo = AuthRepository(tmp_path / "magisteria.sqlite")
    monkeypatch.setattr(application, "auth_repository", repo)
    monkeypatch.setattr(application, "FREE_CUPON_CODES", set())
    admin = authenticated_client()

    page = admin.get("/")
    assert 'id="coupons-button"' in page.text
    created = admin.post("/admin/cupons", json={"cupom": "CATEQUESE", "validade": "semana"})
    assert created.status_code == 200
    assert created.json()["cupom"]["code"] == "CATEQUESE"
    assert created.json()["cupom"]["validity_period"] == "semana"
    assert created.json()["cupom"]["status"] == "ativo"
    coupon_data = created.json()["cupom"]
    assert (datetime.fromisoformat(coupon_data["valid_until"]) - datetime.fromisoformat(coupon_data["created_at"])).days == 7

    user_client = TestClient(application.app)
    create_free_user(user_client, email="promocao@exemplo.com")
    redeemed = user_client.post("/assinatura/cupom", json={"cupom": "catequese"})
    assert redeemed.status_code == 200
    assert redeemed.json()["usuario"]["is_full_access"] is True

    statistics = admin.get("/admin/estatisticas").json()["usuarios"]
    promoted_user = next(user for user in statistics if user["email"] == "promocao@exemplo.com")
    assert promoted_user["coupon_code"] == "CATEQUESE"
    assert promoted_user["can_revoke_coupon"] is True

    revoked = admin.post(
        "/admin/assinatura/revogar-cupom",
        json={"usuario_id": promoted_user["id"]},
    )
    assert revoked.status_code == 200
    assert revoked.json()["usuario"]["is_full_access"] is False
    assert user_client.get("/assinatura").json()["usuario"]["plano"] == "gratuito"

    coupons = admin.get("/admin/cupons").json()["cupons"]
    assert coupons[0]["total_redemptions"] == 1
    assert coupons[0]["active_redemptions"] == 0


def test_managed_coupon_validation_expiration_and_reuse(tmp_path: Path, monkeypatch):
    repo = AuthRepository(tmp_path / "magisteria.sqlite")
    monkeypatch.setattr(application, "auth_repository", repo)
    monkeypatch.setattr(application, "FREE_CUPON_CODES", set())
    admin = authenticated_client()

    assert admin.post("/admin/cupons", json={"cupom": "duas palavras", "validade": "dia"}).status_code == 400
    assert admin.post("/admin/cupons", json={"cupom": "PROMO", "validade": "ano"}).status_code == 400
    assert admin.post("/admin/cupons", json={"cupom": "PROMO", "validade": "dia"}).status_code == 200
    assert admin.post("/admin/cupons", json={"cupom": "promo", "validade": "semana"}).status_code == 400
    monthly = admin.post("/admin/cupons", json={"cupom": "MENSAL", "validade": "mes"}).json()["cupom"]
    monthly_created = datetime.fromisoformat(monthly["created_at"])
    monthly_expiration = datetime.fromisoformat(monthly["valid_until"])
    assert (monthly_expiration.year, monthly_expiration.month) == (
        (monthly_created.year + 1, 1) if monthly_created.month == 12
        else (monthly_created.year, monthly_created.month + 1)
    )

    first_user = TestClient(application.app)
    create_free_user(first_user, email="primeiro@exemplo.com")
    assert first_user.post("/assinatura/cupom", json={"cupom": "PROMO"}).status_code == 200
    first_id = repo.find_user_by_login("primeiro@exemplo.com")["id"]
    admin_id = repo.find_user_by_login("Admin")["id"]
    repo.revoke_coupon_access(first_id, admin_id)
    reused = first_user.post("/assinatura/cupom", json={"cupom": "PROMO"})
    assert reused.status_code == 400
    assert "já foi usado" in reused.json()["detail"]

    with repo._connect() as db:
        db.execute("UPDATE coupons SET valid_until = ? WHERE code = 'PROMO'", ("2020-01-01T00:00:00+00:00",))
    second_user = TestClient(application.app)
    create_free_user(second_user, email="segundo@exemplo.com")
    expired = second_user.post("/assinatura/cupom", json={"cupom": "PROMO"})
    assert expired.status_code == 400
    assert "vencido" in expired.json()["detail"]


def test_payment_with_wrong_amount_is_not_linked(tmp_path: Path, monkeypatch):
    repo = AuthRepository(tmp_path / "magisteria.sqlite")
    monkeypatch.setattr(application, "auth_repository", repo)

    class FakeMercadoPago:
        configured = True
        webhook_signature_configured = True
        price = Decimal("29.90")
        currency = "BRL"

        async def create_subscription(self, user, external_reference):
            self.reference = external_reference
            return {"id": "sub_wrong", "checkout_url": "https://www.mercadopago.com.br/checkout"}

        def validate_webhook_signature(self, *args):
            return True

        async def get_subscription(self, subscription_id):
            return {
                "id": subscription_id,
                "external_reference": self.reference,
                "status": "authorized",
                "auto_recurring": {"transaction_amount": 1, "currency_id": "BRL"},
            }

    provider = FakeMercadoPago()
    monkeypatch.setattr(application, "mercado_pago_service", provider)
    client = TestClient(application.app)
    create_free_user(client, email="valor@exemplo.com")
    mp_user = repo.find_user_by_login("valor@exemplo.com")
    order = repo.create_payment_order(mp_user["id"], "29.90", "BRL", "mercado_pago")
    repo.attach_payment_preference(order["reference"], "sub_wrong")
    provider.reference = order["reference"]

    response = TestClient(application.app).post(
        "/webhooks/mercadopago?data.id=sub_wrong&type=subscription_preapproval",
        headers={"x-signature": "ts=1,v1=fake", "x-request-id": "request-2"},
        json={"type": "subscription_preapproval", "data": {"id": "sub_wrong"}},
    )

    assert response.status_code == 200
    assert repo.find_user_by_login("valor@exemplo.com")["account_type"] == "gratuita"


def test_mercado_pago_webhook_signature_validation():
    service = MercadoPagoService("token", "segredo", Decimal("10.00"), "BRL", "https://app.example")
    timestamp = "1742505638683"
    request_id = "bb56a2f1-6aae-46ac-982e-9dcd3581d08e"
    data_id = "ABC123"
    manifest = f"id:{data_id.lower()};request-id:{request_id};ts:{timestamp};"
    signature = hmac.new(b"segredo", manifest.encode(), hashlib.sha256).hexdigest()

    assert service.validate_webhook_signature(
        f"ts={timestamp},v1={signature}", request_id, data_id
    ) is True
    assert service.validate_webhook_signature(
        f"ts={timestamp},v1={'0' * 64}", request_id, data_id
    ) is False


def test_mercado_pago_subscription_contains_monthly_terms(monkeypatch):
    service = MercadoPagoService(
        "token", "segredo", Decimal("29.90"), "BRL", "https://magisteria.example"
    )
    captured = {}

    async def fake_request(method, path, **kwargs):
        if method == "GET" and path == "/users/me":
            return {"email": "collector@example.com"}
        captured.update(method=method, path=path, payload=kwargs["json"])
        return {
            "id": "sub_456",
            "init_point": "https://www.mercadopago.com.br/checkout/v1/redirect",
        }

    monkeypatch.setattr(service, "_request", fake_request)
    result = asyncio.run(
        service.create_subscription(
            {"full_name": "Maria Silva", "email": "maria@example.com"},
            "mag-reference",
        )
    )

    assert result["id"] == "sub_456"
    assert captured["method"] == "POST"
    assert captured["path"] == "/preapproval"
    assert captured["payload"]["external_reference"] == "mag-reference"
    assert captured["payload"]["payer_email"] == "maria@example.com"
    assert captured["payload"]["auto_recurring"] == {
        "frequency": 1,
        "frequency_type": "months",
        "transaction_amount": 29.9,
        "currency_id": "BRL",
    }
    assert captured["payload"]["status"] == "pending"


def test_mercado_pago_rejects_collector_as_payer(monkeypatch):
    service = MercadoPagoService(
        "token", "segredo", Decimal("14.99"), "BRL", "https://magisteria.example"
    )

    async def fake_request(method, path, **kwargs):
        assert (method, path) == ("GET", "/users/me")
        return {"email": "collector@example.com"}

    monkeypatch.setattr(service, "_request", fake_request)

    with pytest.raises(MercadoPagoError, match="mesma conta") as exc_info:
        asyncio.run(
            service.create_subscription(
                {"full_name": "Conta recebedora", "email": "COLLECTOR@example.com"},
                "mag-reference",
            )
        )
    assert exc_info.value.status_code == 400


def test_asaas_subscription_contains_monthly_terms(monkeypatch):
    service = AsaasService(
        "$aact_hmlg_test-key",
        "webhook-token",
        Decimal("14.99"),
        "https://magisteria.example",
        AsaasService.SANDBOX_URL,
    )
    captured = {}

    async def fake_request(method, path, **kwargs):
        if method == "POST" and path == "/subscriptions":
            captured.update(method=method, path=path, payload=kwargs["json"])
            return {"id": "sub_asaas_123"}
        assert (method, path) == ("GET", "/subscriptions/sub_asaas_123/payments")
        return {
            "data": [
                {
                    "id": "pay_asaas_123",
                    "invoiceUrl": "https://sandbox.asaas.com/i/pay_asaas_123",
                }
            ]
        }

    monkeypatch.setattr(service, "_request", fake_request)
    result = asyncio.run(service.create_subscription("cus_123", "mag-reference"))

    assert result["id"] == "sub_asaas_123"
    assert result["checkout_url"].startswith("https://sandbox.asaas.com/")
    assert captured["payload"]["customer"] == "cus_123"
    assert captured["payload"]["billingType"] == "UNDEFINED"
    assert captured["payload"]["value"] == 14.99
    assert captured["payload"]["cycle"] == "MONTHLY"
    assert captured["payload"]["externalReference"] == "mag-reference"
    assert captured["payload"]["callback"] == {
        "successUrl": "https://magisteria.example/assinatura/retorno",
        "autoRedirect": True,
    }

    service.enable_callback = False
    asyncio.run(service.create_subscription("cus_123", "mag-reference-2"))
    assert "callback" not in captured["payload"]


def test_asaas_rejects_credentials_from_the_wrong_environment():
    common = (
        "webhook-token-with-at-least-thirty-two-characters",
        Decimal("14.99"),
        "https://magisteria.example",
    )
    sandbox_with_production_key = AsaasService(
        "$aact_prod_test-key", *common, AsaasService.SANDBOX_URL
    )
    production_with_sandbox_key = AsaasService(
        "$aact_hmlg_test-key", *common, AsaasService.PRODUCTION_URL
    )
    production = AsaasService(
        "$aact_prod_test-key", *common, AsaasService.PRODUCTION_URL
    )

    assert not sandbox_with_production_key.configured
    assert not production_with_sandbox_key.configured
    assert production.configured


def test_asaas_checkout_and_webhook_control_premium(tmp_path: Path, monkeypatch):
    repo = AuthRepository(tmp_path / "magisteria.sqlite")
    monkeypatch.setattr(application, "auth_repository", repo)

    class FakeAsaas:
        configured = True
        price = Decimal("14.99")
        currency = "BRL"
        status = "PENDING"

        async def get_or_create_customer(self, user, cpf_cnpj):
            assert cpf_cnpj == "24971563792"
            return {"id": "cus_asaas_123"}

        async def create_subscription(self, customer_id, external_reference):
            self.reference = external_reference
            return {
                "id": "sub_asaas_123",
                "payment_id": "pay_asaas_123",
                "checkout_url": "https://sandbox.asaas.com/i/pay_asaas_123",
            }

        async def get_first_subscription_payment(self, subscription_id):
            return {
                "id": "pay_asaas_123",
                "invoiceUrl": "https://sandbox.asaas.com/i/pay_asaas_123",
            }

        async def get_payment(self, payment_id):
            return {
                "id": payment_id,
                "subscription": "sub_asaas_123",
                "externalReference": self.reference,
                "value": 14.99,
                "status": self.status,
            }

        async def get_subscription(self, subscription_id):
            return {
                "id": subscription_id,
                "externalReference": self.reference,
                "nextDueDate": "2026-08-15",
            }

        def validate_webhook_token(self, token):
            return token == "webhook-token"

    provider = FakeAsaas()
    monkeypatch.setattr(application, "asaas_service", provider)
    client = TestClient(application.app)
    create_free_user(client, email="asaas@exemplo.com")

    assert client.post("/assinatura/checkout", json={"cpf_cnpj": "111"}).status_code == 400
    checkout = client.post("/assinatura/checkout", json={"cpf_cnpj": "24971563792"})
    assert checkout.status_code == 200
    order = repo.get_payment_order(checkout.json()["referencia"])
    assert order["provider"] == "asaas"
    assert order["provider_preference_id"] == "sub_asaas_123"
    assert repo.find_user_by_login("asaas@exemplo.com")["account_type"] == "gratuita"

    provider.status = "CONFIRMED"
    event = {
        "id": "evt_asaas_1",
        "event": "PAYMENT_CONFIRMED",
        "payment": {"id": "pay_asaas_123", "status": "FORGED"},
    }
    webhook_client = TestClient(application.app)
    assert webhook_client.post("/webhooks/asaas", json=event).status_code == 401
    approved = webhook_client.post(
        "/webhooks/asaas",
        headers={"asaas-access-token": "webhook-token"},
        json=event,
    )
    assert approved.status_code == 200
    assert approved.json()["status"] == "approved"
    assert repo.find_user_by_login("asaas@exemplo.com")["account_type"] == "completa"

    duplicate = webhook_client.post(
        "/webhooks/asaas",
        headers={"asaas-access-token": "webhook-token"},
        json=event,
    )
    assert duplicate.status_code == 200
    assert "já processada" in duplicate.json()["mensagem"]

    provider.status = "REFUNDED"
    refunded = webhook_client.post(
        "/webhooks/asaas",
        headers={"asaas-access-token": "webhook-token"},
        json={
            "id": "evt_asaas_2",
            "event": "PAYMENT_REFUNDED",
            "payment": {"id": "pay_asaas_123"},
        },
    )
    assert refunded.status_code == 200
    assert repo.find_user_by_login("asaas@exemplo.com")["account_type"] == "gratuita"


def test_source_references_use_document_specific_locators():
    chunks = [
        {"source": "Catecismo da Igreja Católica.pdf", "location": "página 1", "text": "", "score": 1.0, "referencias": ["1210", "1211", "1212"], "categoria": "Catecismo"},
        {"source": "Bíblia Ave Maria - Edição de Estudo.pdf", "location": "São João, página 2", "text": "", "score": 1.0, "referencias": ["Jo 1,2-5"], "categoria": "Bíblia"},
        {"source": "Compêndio Vaticano II.pdf", "location": "página 16", "text": "", "score": 1.0, "referencias": ["Lumen Gentium, n. 1"], "categoria": "Compêndio Vaticano II"},
    ]

    sources = format_sources(chunks)
    abnt = format_abnt_references(chunks)

    assert sources[0]["local"] == "§§ 1210–1212"
    assert sources[1]["local"] == "Jo 1,2-5"
    assert sources[2]["local"] == "Lumen Gentium, n. 1"
    assert "IGREJA CATÓLICA" in abnt
    assert "BÍBLIA. Português." in abnt
    assert "CONCÍLIO VATICANO II" in abnt


def test_abnt_references_skip_weak_auxiliary_sources():
    chunks = [
        {
            "source": "Catecismo da Igreja CatÃ³lica.pdf",
            "location": "pÃ¡gina 10",
            "text": "",
            "score": 0.8,
            "referencias": ["1210"],
            "categoria": "Catecismo",
        },
        {
            "source": "Documento apenas tangencial.pdf",
            "location": "pÃ¡gina 44",
            "text": "",
            "score": 0.1,
            "referencias": [],
            "categoria": "Demais documentos",
        },
    ]

    abnt = format_abnt_references(chunks)

    assert "Catecismo da Igreja" in abnt
    assert "DOCUMENTO APENAS TANGENCIAL" not in abnt


def test_routes_and_closed_base_behavior(monkeypatch):
    client = authenticated_client()
    monkeypatch.setattr(application.vector_store, "search_ordered", lambda *args, **kwargs: [])

    assert client.get("/").status_code == 200
    assert client.get("/health").json()["status"] == "ok"
    status_response = client.get("/status")
    assert status_response.status_code == 200
    assert {"documentos", "trechos", "ultima_atualizacao", "erros"} <= set(status_response.json())
    documents_response = client.get("/documentos")
    assert documents_response.status_code == 200
    assert "documentos" in documents_response.json()
    version_response = client.get("/versao")
    assert version_response.status_code == 200
    assert version_response.json()["versao"] == application.APP_VERSION

    answer_response = client.post("/perguntar", json={"pergunta": "Assunto ausente na base"})
    assert answer_response.status_code == 200
    answer_data = answer_response.json()
    assert answer_data["resposta"] == NOT_FOUND_MESSAGE
    assert answer_data["status_revisao"] == "no_documents"
    assert answer_data["fontes"] == []
    assert answer_data["mensagem_busca"] == ""
    assert answer_data["request_id"]
    stream_response = client.post("/perguntar-stream", json={"pergunta": "Assunto ausente na base"})
    assert stream_response.status_code == 200
    assert '"tipo": "texto"' in stream_response.text
    assert NOT_FOUND_MESSAGE in stream_response.text


def test_followup_prompt_is_inside_answer_below_presentation_buttons():
    client = authenticated_client()

    page = client.get("/").text
    answer_card = page.split('<article class="content-card answer-card">', 1)[1].split("</article>", 1)[0]

    assert page.count('id="followup-panel"') == 1
    assert answer_card.index('id="presentation-module"') < answer_card.index('id="followup-panel"')
    assert answer_card.index('id="create-slides-button"') < answer_card.index('id="followup-slot"')


def test_question_validation():
    client = authenticated_client()
    assert client.post("/perguntar", json={"pergunta": ""}).status_code == 422
    assert application.QuestionRequest(pergunta="Fé").pergunta == "Fé"


def test_rag_diagnostics_redacts_and_persists_trace(tmp_path: Path):
    repository = RAGDiagnosticsRepository(tmp_path / "app.sqlite", debug=True)
    repository.record(
        "req-1",
        "pecado pessoa@example.com 12345678900",
        {
            "query": {"original": "pecado pessoa@example.com 12345678900", "normalized": "pecado", "query_type": "TERM"},
            "candidate_counts": {"lexical_exact": 4},
            "candidates_fused": 4,
            "final_count": 1,
            "reranking": [{"score": 1.2}],
            "selected_chunks": [{"source": "Catecismo.txt"}],
        },
        25,
        "success",
    )

    item = repository.recent(1)[0]
    assert item["query_text"] == ""
    assert item["normalized_query"] == ""
    assert item["trace"]["candidate_counts"]["lexical_exact"] == 4


def test_startup_reindexes_documents(monkeypatch):
    calls = []
    monkeypatch.setattr(application.vector_store, "index_documents", lambda *args: calls.append(True) or {})

    with TestClient(application.app):
        pass

    assert calls == [True]


def test_script_and_slides_routes(monkeypatch):
    client = authenticated_client()
    topics = [{"titulo": "Esperança", "sintese": "Síntese pastoral.", "pontos": ["Primeiro ponto", "Segundo ponto"]}]

    async def outline(*args):
        return topics

    async def plan(*args):
        return {"titulo_curto": "Esperança que transforma", "frase_final": "A esperança renova a vida.", "topicos": topics}

    async def slides(*args):
        return b"pptx"

    monkeypatch.setattr(application.presentation_service, "create_outline", outline)
    monkeypatch.setattr(application.presentation_service, "create_plan", plan)
    monkeypatch.setattr(application.presentation_service, "create_pptx", slides)
    payload = {"titulo": "A esperança cristã", "resposta": "Conteúdo suficientemente longo para gerar materiais."}
    docx = client.post("/criar-roteiro", json=payload)
    pptx = client.post("/criar-slides", json=payload)

    assert docx.status_code == 200
    assert docx.headers["content-type"].startswith("application/vnd.openxmlformats-officedocument.wordprocessingml.document")
    assert pptx.status_code == 200
    assert pptx.content == b"pptx"


def test_generated_docx_and_safe_filename():
    service = PresentationService("", "modelo")
    content = service.create_docx("Fé e esperança", [{"titulo": "Abertura", "sintese": "Uma síntese.", "pontos": ["Um", "Dois"]}])
    assert content.startswith(b"PK")
    assert safe_filename("Fé e esperança", "roteiro.docx") == "fe-e-esperanca-roteiro.docx"


def test_slide_images_are_generated_in_parallel(monkeypatch):
    service = PresentationService("", "modelo", image_concurrency=4)

    async def fake_image(*args):
        await asyncio.sleep(0.05)
        return object()

    monkeypatch.setattr(service, "_generate_image", fake_image)
    topics = [{"titulo": str(index)} for index in range(4)]
    started = time.monotonic()
    images = asyncio.run(service._generate_images("Título", topics))

    assert len(images) == 4
    assert time.monotonic() - started < 0.15


def test_one_failed_image_is_retried_without_losing_the_others(monkeypatch):
    service = PresentationService("", "modelo", image_concurrency=3)
    attempts = {"instavel": 0}

    async def fake_image(_, topic):
        if topic["titulo"] == "instavel":
            attempts["instavel"] += 1
            if attempts["instavel"] == 1:
                raise ValueError("falha temporária")
        return topic["titulo"]

    async def no_wait(*args):
        return None

    monkeypatch.setattr(service, "_generate_image", fake_image)
    monkeypatch.setattr("services.presentation_service.asyncio.sleep", no_wait)
    topics = [{"titulo": "estavel"}, {"titulo": "instavel"}]

    images = asyncio.run(service._generate_images("Título", topics))

    assert images == ["estavel", "instavel"]
    assert attempts["instavel"] == 2
def test_cover_and_closing_images_are_not_hidden_by_black_shapes():
    from io import BytesIO
    from PIL import Image
    from pptx import Presentation
    from pptx.enum.shapes import MSO_SHAPE_TYPE

    image = BytesIO()
    Image.new("RGB", (320, 180), (120, 90, 60)).save(image, "JPEG")
    image.seek(0)
    prs = Presentation()
    PresentationService._add_title_slide(prs, "Título breve", image)
    image.seek(0)
    PresentationService._add_closing_slide(prs, "Síntese final.", image)

    for slide in (prs.slides[-2], prs.slides[-1]):
        assert slide.shapes[0].shape_type == MSO_SHAPE_TYPE.PICTURE
        assert all(shape.shape_type in {MSO_SHAPE_TYPE.PICTURE, MSO_SHAPE_TYPE.TEXT_BOX} for shape in slide.shapes)


def test_user_passwords_must_be_strong_and_can_be_changed(tmp_path: Path):
    repository = AuthRepository(tmp_path / "auth.sqlite")

    ok, message = repository.create_user("Usuario Teste", "usuario@example.com", "fraca1")
    assert not ok
    assert "8 caracteres" in message

    ok, _ = repository.create_user("Usuario Teste", "usuario@example.com", "SenhaForte1")
    assert ok
    user = repository.authenticate("usuario@example.com", "SenhaForte1")

    ok, message = repository.change_password(user["id"], "SenhaForte1", "NovaSenha2")
    assert ok
    assert "sucesso" in message
    assert repository.authenticate("usuario@example.com", "NovaSenha2")


def test_public_database_consolidates_homilies(monkeypatch):
    monkeypatch.setattr(
        application.vector_store,
        "document_names",
        lambda: [
            "Catecismo da Igreja Católica.md",
            "joao-paulo-ii-homilias-txt/1978_homilia.txt",
            "joao-paulo-ii-homilias-txt/1979_homilia.txt",
        ],
    )
    monkeypatch.setattr(application.auth_repository, "inactive_sources", lambda: ())

    names = application.public_document_names()

    assert names.count("Homilias de São João Paulo II") == 1
    assert not any("joao-paulo-ii-homilias-txt" in name for name in names)


def test_password_protects_application_and_health_is_public():
    client = TestClient(application.app)
    assert client.get("/health").status_code == 200
    assert client.get("/", follow_redirects=False).headers["location"] == "/login"
    admin = authenticated_client()
    assert admin.get("/").status_code == 200


def test_admin_can_clear_document_base(monkeypatch, tmp_path: Path):
    documents = tmp_path / "Documentos"
    documents.mkdir()
    (documents / "antigo.pdf").write_bytes(b"%PDF-1.4")
    (documents / "antigo.md").write_text("texto antigo", encoding="utf-8")
    (documents / "manter.png").write_bytes(b"imagem")

    async def fake_indexing():
        return {"documentos": 0, "trechos": 0}

    monkeypatch.setattr(application.settings, "DOCUMENTS_DIR", documents)
    monkeypatch.setattr(application, "perform_indexing", fake_indexing)

    client = authenticated_client()
    response = client.post("/admin/base-documental/limpar")

    assert response.status_code == 200
    assert sorted(response.json()["removidos"]) == ["antigo.md", "antigo.pdf"]
    assert not (documents / "antigo.pdf").exists()
    assert not (documents / "antigo.md").exists()
    assert (documents / "manter.png").exists()


def test_admin_can_upload_document_chunks(monkeypatch, tmp_path: Path):
    documents = tmp_path / "Documentos"
    documents.mkdir()
    monkeypatch.setattr(application.settings, "DOCUMENTS_DIR", documents)

    client = authenticated_client()
    response = client.post(
        "/admin/upload-chunk",
        content=b"parte 1 ",
        headers={
            "X-Path": "subpasta/documento.txt",
            "X-Offset": "0",
            "X-Complete": "0",
        },
    )
    second = client.post(
        "/admin/upload-chunk",
        content=b"parte 2",
        headers={
            "X-Path": "subpasta/documento.txt",
            "X-Offset": "8",
            "X-Complete": "1",
        },
    )

    assert response.status_code == 200
    assert second.status_code == 200
    assert second.json()["arquivo"] == "subpasta/documento.txt"
    assert (documents / "subpasta" / "documento.txt").read_text(encoding="utf-8") == "parte 1 parte 2"


def test_localized_no_document_messages_do_not_require_api_key():
    service = AnswerService("", "modelo")

    async def call(language):
        return await service.answer_with_review("Question", [], language=language)

    english = asyncio.run(call("en"))
    spanish = asyncio.run(call("es"))

    assert english["resposta"] == answer_message("no_documents", "en")
    assert spanish["resposta"] == answer_message("no_documents", "es")


def test_foreign_answer_prompt_requires_only_the_selected_language():
    service = AnswerService("chave", "modelo")
    chunks = [{"source": "Catecismo.txt", "location": "página 1", "text": "A dignidade é inerente à pessoa.", "score": 1.0}]

    english = service._request_arguments("What is human dignity?", chunks, [], language="en")
    spanish = service._request_arguments("¿Qué es la dignidad humana?", chunks, [], language="es")

    assert "exclusively in clear, natural international English" in english["instructions"]
    assert "exclusivamente en español internacional" in spanish["instructions"]
    assert "em português brasileiro contemporâneo" not in english["instructions"]
    assert "em português brasileiro contemporâneo" not in spanish["instructions"]


def test_query_translation_returns_portuguese_search_text(monkeypatch):
    class FakeResponses:
        def __init__(self):
            self.calls = []

        async def create(self, **kwargs):
            self.calls.append(kwargs)
            return SimpleNamespace(output_text="O que é a dignidade da pessoa humana?")

    fake_responses = FakeResponses()
    monkeypatch.setattr(
        "services.answer_service.AsyncOpenAI",
        lambda api_key: SimpleNamespace(responses=fake_responses),
    )
    service = AnswerService("chave", "modelo")

    translated = asyncio.run(service.translate_query_to_portuguese("What is human dignity?", "en"))

    assert translated == "O que é a dignidade da pessoa humana?"
    assert fake_responses.calls[0]["input"] == "What is human dignity?"
    assert "uso exclusivo em uma busca documental" in fake_responses.calls[0]["instructions"]


def test_foreign_query_is_translated_before_retrieval_and_answered_in_selected_language(tmp_path: Path, monkeypatch):
    repository = AuthRepository(tmp_path / "magisteria.sqlite")
    diagnostics_repository = RAGDiagnosticsRepository(tmp_path / "magisteria.sqlite", False)
    monkeypatch.setattr(application, "auth_repository", repository)
    monkeypatch.setattr(application, "rag_diagnostics", diagnostics_repository)
    calls = {}

    class FakeAnswerService:
        api_key = "configured"

        async def translate_query_to_portuguese(self, query, source_language):
            calls["translation"] = (query, source_language)
            return "dignidade da pessoa humana"

        async def answer_with_review(self, question, chunks, history, style_chunks, language):
            calls["answer"] = (question, language)
            return {"resposta": "Human dignity is inherent to every person.", "status_revisao": "approve", "motivo_revisao": ""}

    def fake_retrieval(payload, query_override=None):
        calls["retrieval"] = query_override
        return ([{
            "source": "Catecismo.txt", "location": "página 1", "text": "A dignidade é inerente.",
            "score": 1.0, "ordem": 1, "categoria": "Documento",
        }], {"query": {"query_type": "question"}, "reranking": [{"score_normalized": 1.0, "strategies": ["lexical_exact"]}]})

    monkeypatch.setattr(application, "answer_service", FakeAnswerService())
    monkeypatch.setattr(application, "ordered_chunks_with_diagnostics", fake_retrieval)
    monkeypatch.setattr(application, "homily_style_chunks", lambda payload, query_override=None: [])
    client = TestClient(application.app)
    create_free_user(client, email="polyglot@example.com")

    response = client.post("/perguntar", json={"pergunta": "What is human dignity?", "idioma": "en"})

    assert response.status_code == 200
    assert calls["translation"] == ("What is human dignity?", "en")
    assert calls["retrieval"] == "dignidade da pessoa humana"
    assert calls["answer"] == ("What is human dignity?", "en")
    assert response.json()["resposta"] == "Human dignity is inherent to every person."


def test_admin_bootstrap_requires_an_explicit_strong_secret(tmp_path: Path):
    repository = AuthRepository(tmp_path / "without-admin.sqlite")
    assert repository.find_user_by_login("Admin") is None

    repository = AuthRepository(
        tmp_path / "with-admin.sqlite",
        admin_bootstrap_password="AdminTest3510",
    )
    admin = repository.authenticate("Admin", "AdminTest3510")
    assert admin is not None
    assert admin["role"] == "admin"
    assert repository.authenticate("Admin", "3510") is None


def test_password_change_revokes_every_existing_session(tmp_path: Path):
    repository = AuthRepository(tmp_path / "sessions.sqlite")
    ok, _ = repository.create_user("Usuario Teste", "session@example.com", "SenhaForte1")
    assert ok
    user = repository.authenticate("session@example.com", "SenhaForte1")
    first = repository.create_session(user["id"])
    second = repository.create_session(user["id"])

    ok, _ = repository.change_password(user["id"], "SenhaForte1", "NovaSenha2")

    assert ok
    assert repository.get_user_by_session(first) is None
    assert repository.get_user_by_session(second) is None
    assert repository.authenticate("session@example.com", "NovaSenha2") is not None


def test_usage_reservations_are_atomic_under_concurrency(tmp_path: Path):
    repository = AuthRepository(tmp_path / "quota.sqlite")
    ok, _ = repository.create_user("Usuario Teste", "quota@example.com", "SenhaForte1")
    assert ok
    user_id = repository.find_user_by_login("quota@example.com")["id"]

    with ThreadPoolExecutor(max_workers=8) as executor:
        query_results = list(executor.map(lambda _: repository.reserve_usage(user_id, "query"), range(8)))
        script_results = list(executor.map(lambda _: repository.reserve_usage(user_id, "script"), range(4)))
        slide_results = list(executor.map(lambda _: repository.reserve_usage(user_id, "presentation"), range(4)))

    assert sum(allowed for allowed, _ in query_results) == 3
    assert sum(allowed for allowed, _ in script_results) == 1
    assert sum(allowed for allowed, _ in slide_results) == 1


def test_ordinary_user_cannot_trigger_legacy_reindex(tmp_path: Path, monkeypatch):
    repository = AuthRepository(tmp_path / "reindex.sqlite")
    monkeypatch.setattr(application, "auth_repository", repository)
    client = TestClient(application.app)
    create_free_user(client, email="reindex@example.com")

    response = client.post("/indexar")

    assert response.status_code == 403


def test_mercado_pago_webhook_fails_closed_without_secret():
    service = MercadoPagoService("token", "", Decimal("10.00"), "BRL", "https://app.example")
    assert service.validate_webhook_signature("ts=1,v1=hash", "request-1", "payment-1") is False


def test_slide_fallback_does_not_block_the_event_loop(monkeypatch):
    from io import BytesIO
    from PIL import Image
    from pptx import Presentation

    service = PresentationService("", "modelo")
    Presentation()

    def image_bytes():
        output = BytesIO()
        Image.new("RGB", (64, 64), (30, 40, 50)).save(output, "JPEG")
        output.seek(0)
        return output

    async def generated_image(*args, **kwargs):
        return image_bytes()

    def slow_fallback(*args, **kwargs):
        time.sleep(0.08)
        return image_bytes()

    monkeypatch.setattr(service, "_generate_image_with_retry", generated_image)
    monkeypatch.setattr(service, "_fallback_image", slow_fallback)
    image_bytes()

    async def scenario():
        task = asyncio.create_task(
            service.create_pptx(
                "Tema",
                [{"titulo": "Topico", "sintese": "Sintese", "pontos": ["Ponto"]}],
            )
        )
        started = time.monotonic()
        await asyncio.sleep(0.02)
        heartbeat_elapsed = time.monotonic() - started
        await task
        return heartbeat_elapsed

    assert asyncio.run(scenario()) < 0.06


def test_mobile_tokens_rotate_and_reuse_revokes_the_family(tmp_path: Path, monkeypatch):
    repository = AuthRepository(tmp_path / "mobile-auth.sqlite")
    monkeypatch.setattr(application, "auth_repository", repository)
    assert repository.create_user("Usuario Mobile", "mobile@example.com", "SenhaForte1")[0]
    client = TestClient(application.app)

    login = client.post(
        "/api/v1/mobile/auth/login",
        json={"email": "mobile@example.com", "password": "SenhaForte1"},
    )
    assert login.status_code == 200
    initial = login.json()
    me = client.get(
        "/api/v1/mobile/me",
        headers={"Authorization": f"Bearer {initial['access_token']}"},
    )
    assert me.status_code == 200
    assert me.json()["user"]["email"] == "mobile@example.com"

    rotated = client.post(
        "/api/v1/mobile/auth/refresh",
        json={"refresh_token": initial["refresh_token"]},
    )
    assert rotated.status_code == 200
    current = rotated.json()
    reuse = client.post(
        "/api/v1/mobile/auth/refresh",
        json={"refresh_token": initial["refresh_token"]},
    )
    assert reuse.status_code == 401
    assert client.get(
        "/api/v1/mobile/me",
        headers={"Authorization": f"Bearer {current['access_token']}"},
    ).status_code == 401


def test_mobile_account_deletion_reauthenticates_and_cascades(tmp_path: Path, monkeypatch):
    repository = AuthRepository(tmp_path / "mobile-delete.sqlite")
    monkeypatch.setattr(application, "auth_repository", repository)
    assert repository.create_user("Usuario Mobile", "delete@example.com", "SenhaForte1")[0]
    client = TestClient(application.app)
    login = client.post(
        "/api/v1/mobile/auth/login",
        json={"email": "delete@example.com", "password": "SenhaForte1"},
    ).json()
    headers = {"Authorization": f"Bearer {login['access_token']}"}

    denied = client.request(
        "DELETE",
        "/api/v1/mobile/account",
        headers=headers,
        json={"password": "errada", "confirmation": "EXCLUIR"},
    )
    assert denied.status_code == 401
    deleted = client.request(
        "DELETE",
        "/api/v1/mobile/account",
        headers=headers,
        json={"password": "SenhaForte1", "confirmation": "EXCLUIR"},
    )
    assert deleted.status_code == 200
    assert repository.find_user_by_login("delete@example.com") is None
    assert client.get("/api/v1/mobile/me", headers=headers).status_code == 401
    with repository._connect() as db:
        assert db.execute("SELECT COUNT(*) FROM account_deletion_audit").fetchone()[0] == 1


def test_security_headers_and_mobile_cors_are_present():
    client = TestClient(application.app)
    response = client.get("/health", headers={"Origin": "capacitor://localhost"})
    assert response.headers["x-content-type-options"] == "nosniff"
    assert response.headers["x-frame-options"] == "DENY"
    assert "frame-ancestors 'none'" in response.headers["content-security-policy"]
    assert response.headers["access-control-allow-origin"] == "capacitor://localhost"


def test_subscription_abstraction_normalizes_web_android_ios_and_free():
    service = SubscriptionService("google.premium", "apple.premium")
    user = {
        "role": "user",
        "account_type": "completa",
        "subscription_status": "ativa",
        "subscription_renews_at": "2026-08-01",
    }
    assert service.snapshot(user, "asaas").source == "web"
    assert service.snapshot(user, "google_play").product_id == "google.premium"
    assert service.snapshot(user, "apple").source == "ios"
    free = {**user, "account_type": "gratuita", "subscription_status": "inativa"}
    assert service.snapshot(free).source == "free"
    assert service.snapshot(free).is_full_access is False
