from __future__ import annotations

import base64
import asyncio
import hashlib
from io import BytesIO
import json
import logging
import re
import unicodedata

from docx import Document
from docx.enum.text import WD_ALIGN_PARAGRAPH
from docx.shared import Inches, Pt, RGBColor
import httpx
from openai import AsyncOpenAI
from openai import APIConnectionError, APIStatusError, RateLimitError
from PIL import Image, ImageDraw


MIN_PRESENTATION_TOPICS = 10
MAX_PRESENTATION_TOPICS = 14
logger = logging.getLogger(__name__)


class PresentationService:
    def __init__(
        self,
        api_key: str,
        text_model: str,
        image_model: str = "gpt-image-1",
        image_concurrency: int = 4,
        image_quality: str = "low",
    ):
        self.client = AsyncOpenAI(api_key=api_key) if api_key else None
        self.text_model = text_model
        self.image_model = image_model
        self.image_concurrency = max(1, min(image_concurrency, 6))
        self.image_quality = image_quality

    async def create_plan(self, title: str, answer: str) -> dict:
        if not self.client:
            raise RuntimeError("A chave OPENAI_API_KEY ainda não foi configurada no arquivo .env.")
        response = await self.client.responses.create(
            model=self.text_model,
            instructions=(
                "Organize exclusivamente o conteúdo fornecido em 10 a 14 tópicos para uma pregação ou palestra. "
                "Não acrescente fatos, citações ou referências externas. Cada tópico deve ter título curto, "
                "síntese de até 38 palavras e 2 a 3 pontos de desenvolvimento. Divida ideias amplas em etapas "
                "menores para que a apresentação tenha ritmo, progressão e mais slides. Crie também um título "
                "de capa marcante, coerente e com no máximo 7 palavras, e uma frase final de até 24 palavras que "
                "resuma a mensagem central. Responda somente em JSON válido."
            ),
            input=f"TÍTULO: {title}\n\nCONTEÚDO: {answer}",
            text={"format": {"type": "json_schema", "name": "roteiro", "strict": True, "schema": {
                "type": "object", "properties": {
                "titulo_curto": {"type": "string"},
                "frase_final": {"type": "string"},
                "topicos": {"type": "array", "minItems": MIN_PRESENTATION_TOPICS, "maxItems": MAX_PRESENTATION_TOPICS,
                "items": {"type": "object", "properties": {
                    "titulo": {"type": "string"}, "sintese": {"type": "string"},
                    "pontos": {"type": "array", "minItems": 2, "maxItems": 3, "items": {"type": "string"}}
                }, "required": ["titulo", "sintese", "pontos"], "additionalProperties": False}}},
                "required": ["titulo_curto", "frase_final", "topicos"], "additionalProperties": False}}},
            max_output_tokens=3600,
        )
        return json.loads(response.output_text)

    async def create_outline(self, title: str, answer: str) -> list[dict]:
        return (await self.create_plan(title, answer))["topicos"]

    def create_docx(self, title: str, topics: list[dict]) -> bytes:
        doc = Document()
        section = doc.sections[0]
        section.top_margin = section.bottom_margin = Inches(0.8)
        section.left_margin = section.right_margin = Inches(0.9)
        doc.styles["Normal"].font.name = "Aptos"
        doc.styles["Normal"].font.size = Pt(11)
        doc.styles["Heading 1"].font.name = "Aptos Display"
        doc.styles["Heading 1"].font.size = Pt(16)
        doc.styles["Heading 1"].font.color.rgb = RGBColor(91, 55, 28)
        kicker = doc.add_paragraph()
        kicker.alignment = WD_ALIGN_PARAGRAPH.CENTER
        run = kicker.add_run("MAGISTERIA · ROTEIRO PARA PREGAÇÃO OU PALESTRA")
        run.bold = True
        run.font.size = Pt(9)
        run.font.color.rgb = RGBColor(174, 124, 55)
        heading = doc.add_paragraph()
        heading.alignment = WD_ALIGN_PARAGRAPH.CENTER
        title_run = heading.add_run(title)
        title_run.bold = True
        title_run.font.name = "Aptos Display"
        title_run.font.size = Pt(25)
        title_run.font.color.rgb = RGBColor(54, 38, 29)
        doc.add_paragraph("Tópicos organizados a partir da pesquisa realizada no MAGISTERIA.").alignment = WD_ALIGN_PARAGRAPH.CENTER
        for number, topic in enumerate(topics, 1):
            doc.add_heading(f"{number}. {topic['titulo']}", level=1)
            doc.add_paragraph(topic["sintese"])
            for point in topic["pontos"]:
                doc.add_paragraph(point, style="List Bullet")
        footer = section.footer.paragraphs[0]
        footer.alignment = WD_ALIGN_PARAGRAPH.CENTER
        footer.add_run("MAGISTERIA · Conteúdo fundamentado na pesquisa selecionada").font.size = Pt(8)
        output = BytesIO()
        doc.save(output)
        return output.getvalue()

    async def create_pptx(
        self,
        title: str,
        topics: list[dict],
        short_title: str | None = None,
        closing_phrase: str | None = None,
    ) -> bytes:
        from pptx import Presentation
        from pptx.util import Inches as PptInches

        short_title = short_title or title
        closing_phrase = closing_phrase or "A mensagem acolhida com fé transforma a vida e renova a esperança."
        visual_brief = self._visual_brief(title, topics, short_title, closing_phrase)
        images = await self._generate_images(title, visual_brief)
        prs = Presentation()
        prs.slide_width, prs.slide_height = PptInches(13.333), PptInches(7.5)
        self._add_title_slide(prs, short_title, images[0])
        for number, (topic, image) in enumerate(zip(topics, images[1:-1]), 1):
            self._add_topic_slide(prs, number, topic, image)
        self._add_closing_slide(prs, closing_phrase, images[-1])
        output = BytesIO()
        prs.save(output)
        return output.getvalue()

    def _visual_brief(self, title: str, topics: list[dict], short_title: str, closing_phrase: str) -> list[dict]:
        total = len(topics) + 2
        brief = [
            {"titulo": short_title, "sintese": title, "visual_role": "cover"},
            *topics,
            {"titulo": "Encerramento", "sintese": closing_phrase, "visual_role": "closing"},
        ]
        return [
            {**item, "visual_index": index, "visual_total": total, "visual_signature": self._visual_signature(index)}
            for index, item in enumerate(brief, 1)
        ]

    async def _generate_images(self, title: str, topics: list[dict]) -> list[BytesIO]:
        semaphore = asyncio.Semaphore(self.image_concurrency)

        async def generate(topic: dict) -> BytesIO:
            async with semaphore:
                return await self._generate_image_with_retry(title, topic)

        results = await asyncio.gather(*(generate(topic) for topic in topics), return_exceptions=True)
        # Um limite momentâneo em uma imagem não deve descartar todas as demais.
        for index, result in enumerate(results):
            if isinstance(result, Exception):
                await asyncio.sleep(1)
                try:
                    results[index] = await self._generate_image_with_retry(title, topics[index], attempts=3)
                except Exception as exc:
                    logger.warning("Falha ao gerar imagem do slide %s; usando imagem reserva.", index + 1, exc_info=exc)
                    results[index] = self._fallback_image(topics[index])
        seen: set[str] = set()
        for index, result in enumerate(results):
            digest = self._image_digest(result)
            if digest in seen:
                unique_topic = {
                    **topics[index],
                    "visual_signature": f"{topics[index].get('visual_signature', '')}; composição alternativa sem repetir imagens anteriores",
                }
                try:
                    results[index] = await self._generate_image_with_retry(title, unique_topic, attempts=2)
                except Exception as exc:
                    logger.warning("Falha ao regenerar imagem repetida do slide %s; usando imagem reserva.", index + 1, exc_info=exc)
                    results[index] = self._fallback_image(unique_topic)
                digest = self._image_digest(results[index])
            seen.add(digest)
        return list(results)

    async def _generate_image_with_retry(self, title: str, topic: dict, attempts: int = 2) -> BytesIO:
        for attempt in range(attempts):
            try:
                return await self._generate_image(title, topic)
            except (RateLimitError, APIConnectionError) as exc:
                if attempt + 1 == attempts:
                    raise RuntimeError(
                        "O serviço de imagens está temporariamente ocupado. Aguarde alguns instantes e tente novamente."
                    ) from exc
                await asyncio.sleep(1.5 * (attempt + 1))
            except APIStatusError as exc:
                if exc.status_code >= 500 and attempt + 1 < attempts:
                    await asyncio.sleep(1.5 * (attempt + 1))
                    continue
                raise
        raise RuntimeError("Não foi possível concluir a geração das imagens.")

    async def _generate_image(self, title: str, topic: dict) -> BytesIO:
        role = topic.get("visual_role", "topic")
        direction = {
            "cover": (
                "Imagem de capa impactante, majestosa e emocional, com forte profundidade, luz dramática e um "
                "ponto focal simbólico. Reserve área escura e limpa no centro para sobrepor um título."
            ),
            "closing": (
                "Imagem final luminosa, serena e esperançosa, com atmosfera de comunhão, gratidão e bênção. "
                "Reserve espaço central limpo para uma frase de encerramento."
            ),
            "topic": "Imagem narrativa coerente com o tópico, com foco visual claro e composição elegante.",
        }[role]
        prompt = (
            "Ilustração editorial cinematográfica, reverente e acolhedora, sem texto, letras ou logotipos, "
            "para apresentação católica contemporânea. Evite retratar Deus literalmente. "
            "Cada slide desta apresentação deve ter imagem inédita; não reutilize personagens, enquadramento, "
            "cenário ou composição dos demais slides. "
            f"Imagem {topic.get('visual_index', 1)} de {topic.get('visual_total', 1)}. "
            f"Identidade visual exclusiva: {topic.get('visual_signature', '')}. "
            f"{direction} Tema geral: {title}. Tópico: {topic['titulo']}. Ideia: {topic['sintese']}"
        )
        result = await self.client.images.generate(**self._image_request_args(prompt))
        image = result.data[0]
        if getattr(image, "b64_json", None):
            return BytesIO(base64.b64decode(image.b64_json))
        if getattr(image, "url", None):
            async with httpx.AsyncClient(timeout=60) as client:
                response = await client.get(image.url)
                response.raise_for_status()
                return BytesIO(response.content)
        raise RuntimeError("O serviço de imagens não retornou um arquivo utilizável.")

    def _image_request_args(self, prompt: str) -> dict:
        args = {
            "model": self.image_model,
            "prompt": prompt,
            "size": "1024x1024",
        }
        if self.image_model.startswith("dall-e"):
            args["quality"] = "hd" if self.image_quality == "high" else "standard"
            args["response_format"] = "b64_json"
        else:
            args["quality"] = self.image_quality
            args["output_format"] = "jpeg"
        return args

    @staticmethod
    def _fallback_image(topic: dict) -> BytesIO:
        index = int(topic.get("visual_index", 1))
        palettes = [
            ((40, 28, 22), (165, 111, 45), (250, 238, 209)),
            ((19, 45, 37), (92, 134, 95), (238, 226, 190)),
            ((30, 36, 58), (119, 91, 150), (244, 231, 211)),
            ((55, 33, 31), (150, 67, 48), (251, 230, 196)),
            ((31, 45, 62), (62, 111, 145), (235, 226, 203)),
            ((47, 39, 22), (154, 127, 59), (248, 238, 214)),
        ]
        dark, mid, light = palettes[(index - 1) % len(palettes)]
        image = Image.new("RGB", (1024, 1024), dark)
        pixels = image.load()
        for y in range(1024):
            ratio = y / 1023
            base = tuple(int(dark[channel] * (1 - ratio) + mid[channel] * ratio) for channel in range(3))
            for x in range(1024):
                glow = max(0, 1 - (((x - 670) ** 2 + (y - 360) ** 2) ** 0.5 / 760))
                pixels[x, y] = tuple(
                    min(255, int(base[channel] * (1 - glow * 0.5) + light[channel] * glow * 0.5))
                    for channel in range(3)
                )
        draw = ImageDraw.Draw(image, "RGBA")
        for offset in range(0, 1024, 96):
            alpha = 32 + ((offset + index * 17) % 48)
            draw.line([(offset - 260, 1024), (offset + 360, 0)], fill=(*light, alpha), width=10)
        draw.ellipse((650, 160, 930, 440), fill=(*light, 46))
        draw.ellipse((700, 210, 880, 390), fill=(*light, 54))
        output = BytesIO()
        image.save(output, "JPEG", quality=86)
        output.seek(0)
        return output

    @staticmethod
    def _visual_signature(index: int) -> str:
        signatures = [
            "plano aberto com arquitetura sacra e luz atravessando o ambiente",
            "detalhe simbólico de mãos, livro e vela em atmosfera de recolhimento",
            "caminho externo ao amanhecer com profundidade e horizonte luminoso",
            "comunidade reunida em silhuetas discretas, clima de escuta e unidade",
            "mesa simples com Bíblia aberta, textura artesanal e luz lateral quente",
            "vitrais abstratos projetando cores suaves sem texto nem figuras literais",
            "paisagem serena com cruz distante, nuvens dramáticas e luz de esperança",
            "interior contemporâneo minimalista com foco em símbolo sacramental",
            "peregrinação em estrada de pedra, movimento sutil e perspectiva baixa",
            "luz dourada sobre detalhes de madeira, tecido e objeto litúrgico discreto",
            "jardim silencioso após a chuva, atmosfera de renovação e contemplação",
            "assembleia em penumbra com feixe de luz central e composição cinematográfica",
            "porta aberta para ambiente iluminado, símbolo de passagem e conversão",
            "céu amplo com raios de sol e elementos naturais em composição reverente",
            "capela pequena vista em diagonal, profundidade elegante e sombras suaves",
            "detalhe de passos no chão, sinal de missão e continuidade",
        ]
        return signatures[(index - 1) % len(signatures)]

    @staticmethod
    def _image_digest(image: BytesIO) -> str:
        if not hasattr(image, "tell") or not hasattr(image, "seek") or not hasattr(image, "read"):
            return hashlib.sha256(repr(image).encode("utf-8")).hexdigest()
        position = image.tell()
        image.seek(0)
        digest = hashlib.sha256(image.read()).hexdigest()
        image.seek(position)
        return digest

    @staticmethod
    def _add_title_slide(prs, title: str, image: BytesIO) -> None:
        from pptx.dml.color import RGBColor as PptRGBColor
        from pptx.enum.text import PP_ALIGN
        from pptx.util import Inches as PptInches, Pt as PptPt

        slide = prs.slides.add_slide(prs.slide_layouts[6])
        slide.shapes.add_picture(image, 0, 0, width=prs.slide_width, height=prs.slide_height)
        frame = slide.shapes.add_textbox(PptInches(1.25), PptInches(2.05), PptInches(10.83), PptInches(3.5)).text_frame
        p = frame.paragraphs[0]
        p.text, p.alignment = title, PP_ALIGN.CENTER
        p.font.name, p.font.size, p.font.bold = "Aptos Display", PptPt(50), True
        p.font.color.rgb = PptRGBColor(255, 248, 235)
        sub = frame.add_paragraph()
        sub.text, sub.alignment = "Gerado com curadoria via MagisterIA", PP_ALIGN.CENTER
        sub.space_before = PptPt(22)
        sub.font.size, sub.font.color.rgb = PptPt(18), PptRGBColor(255, 238, 207)

    @staticmethod
    def _add_topic_slide(prs, number: int, topic: dict, image: BytesIO) -> None:
        from pptx.dml.color import RGBColor as PptRGBColor
        from pptx.util import Inches as PptInches, Pt as PptPt

        slide = prs.slides.add_slide(prs.slide_layouts[6])
        slide.background.fill.solid()
        slide.background.fill.fore_color.rgb = PptRGBColor(250, 246, 238)
        slide.shapes.add_picture(image, PptInches(6.55), 0, width=PptInches(6.783), height=PptInches(7.5))
        frame = slide.shapes.add_textbox(PptInches(0.7), PptInches(0.55), PptInches(5.35), PptInches(6.3)).text_frame
        frame.word_wrap = True
        p = frame.paragraphs[0]
        p.text = f"{number:02d}  {topic['titulo']}"
        p.font.name, p.font.size, p.font.bold = "Aptos Display", PptPt(35), True
        p.font.color.rgb = PptRGBColor(70, 44, 30)
        lead = frame.add_paragraph()
        lead.text, lead.space_before, lead.space_after = topic["sintese"], PptPt(18), PptPt(14)
        lead.font.size, lead.font.color.rgb = PptPt(20), PptRGBColor(91, 71, 58)
        for point in topic["pontos"]:
            bullet = frame.add_paragraph()
            bullet.text, bullet.font.size, bullet.space_after = f"• {point}", PptPt(17), PptPt(8)
            bullet.font.color.rgb = PptRGBColor(54, 45, 40)

    @staticmethod
    def _add_closing_slide(prs, phrase: str, image: BytesIO) -> None:
        from pptx.dml.color import RGBColor as PptRGBColor
        from pptx.enum.text import PP_ALIGN
        from pptx.util import Inches as PptInches, Pt as PptPt

        slide = prs.slides.add_slide(prs.slide_layouts[6])
        slide.shapes.add_picture(image, 0, 0, width=prs.slide_width, height=prs.slide_height)
        frame = slide.shapes.add_textbox(PptInches(1.3), PptInches(1.45), PptInches(10.73), PptInches(4.9)).text_frame
        frame.word_wrap = True
        p = frame.paragraphs[0]
        p.text, p.alignment = phrase, PP_ALIGN.CENTER
        p.font.name, p.font.size, p.font.bold = "Aptos Display", PptPt(34), True
        p.font.color.rgb = PptRGBColor(255, 250, 239)
        thanks = frame.add_paragraph()
        thanks.text, thanks.alignment = "Obrigado pela atenção!", PP_ALIGN.CENTER
        thanks.space_before = PptPt(28)
        thanks.font.size, thanks.font.bold = PptPt(24), True
        thanks.font.color.rgb = PptRGBColor(239, 202, 132)
        blessing = frame.add_paragraph()
        blessing.text, blessing.alignment = "Que a bênção de Deus esteja com todos!", PP_ALIGN.CENTER
        blessing.space_before = PptPt(12)
        blessing.font.size = PptPt(22)
        blessing.font.color.rgb = PptRGBColor(255, 248, 235)


def safe_filename(title: str, suffix: str) -> str:
    value = unicodedata.normalize("NFKD", title).encode("ascii", "ignore").decode("ascii")
    slug = re.sub(r"[^a-zA-Z0-9]+", "-", value).strip("-").lower()[:60] or "magisteria"
    return f"{slug}-{suffix}"
